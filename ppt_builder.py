import logging
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

STYLES = {
    "academic": {
        "name": "学术期刊风",
        "slide_bg": "FFFFFF",
        "alt_bg": "F7F9FC",
        "accent": "1A5276",
        "accent2": "2E86C1",
        "text": "1C2833",
        "muted": "5D6D7E",
        "light_bg": "EBF5FB",
        "highlight_bg": "FEF9E7",
        "title_size": Pt(32),
        "subtitle_size": Pt(18),
        "body_size": Pt(14),
        "small_size": Pt(11),
        "font": "Microsoft YaHei",
    },
    "tech_blueprint": {
        "name": "科技蓝图风",
        "slide_bg": "F0F4F8",
        "alt_bg": "E8EEF5",
        "accent": "1B4F72",
        "accent2": "00A8CC",
        "text": "1C2E3F",
        "muted": "5D7B93",
        "light_bg": "D6E4F0",
        "highlight_bg": "E8F8F5",
        "title_size": Pt(30),
        "subtitle_size": Pt(16),
        "body_size": Pt(13),
        "small_size": Pt(10),
        "font": "Segoe UI",
    },
    "experiment": {
        "name": "实验报告风",
        "slide_bg": "FAFAFA",
        "alt_bg": "F2F3F4",
        "accent": "1E8449",
        "accent2": "27AE60",
        "text": "1C2833",
        "muted": "616A6B",
        "light_bg": "EAFAF1",
        "highlight_bg": "FDEBD0",
        "title_size": Pt(28),
        "subtitle_size": Pt(16),
        "body_size": Pt(13),
        "small_size": Pt(10),
        "font": "Arial",
    },
    "chinese_meeting": {
        "name": "中文组会风",
        "slide_bg": "FFFBF0",
        "alt_bg": "F8F4EA",
        "accent": "8B6914",
        "accent2": "D4A017",
        "text": "3D2B1F",
        "muted": "8B7D6B",
        "light_bg": "FDF2E9",
        "highlight_bg": "FEF5E7",
        "title_size": Pt(34),
        "subtitle_size": Pt(18),
        "body_size": Pt(15),
        "small_size": Pt(11),
        "font": "Microsoft YaHei",
    },
}


def _hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _add_bg(slide, color_hex):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(color_hex)


def _add_textbox(slide, left, top, width, height, text, font_size, color_hex, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = font_size
    p.font.color.rgb = _hex_to_rgb(color_hex)
    p.font.bold = bold
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return txBox


def _add_bullet_textbox(slide, left, top, width, height, lines, font_size, color_hex, font_name=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = str(line)[:200]
        p.font.size = font_size
        p.font.color.rgb = _hex_to_rgb(color_hex)
        p.space_after = Pt(6)
        if font_name:
            p.font.name = font_name
    return txBox


def _add_shape(slide, shape_type, left, top, width, height, fill_hex=None, line_hex=None):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    else:
        shape.fill.background()
    if line_hex:
        shape.line.color.rgb = _hex_to_rgb(line_hex)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_hex, text="", font_size=Pt(12), font_color="FFFFFF", font_name=None):
    shape = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_hex)
    shape.text_frame.word_wrap = True
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].text = str(text)[:80]
    shape.text_frame.paragraphs[0].font.size = font_size
    shape.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(font_color)
    if font_name:
        shape.text_frame.paragraphs[0].font.name = font_name
    shape.text_frame.paragraphs[0].font.bold = True
    return shape


def _add_arrow(slide, left, top, width, height, color_hex):
    shape = _add_shape(slide, MSO_SHAPE.RIGHT_ARROW, left, top, width, height, color_hex)
    return shape


class PPTBuilder:
    def __init__(self, style="academic"):
        self.style_config = STYLES.get(style, STYLES["academic"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self._slide_layout = self.prs.slide_layouts[6]

    def _slide_number(self, idx, total):
        return f"{idx}/{total}"

    def _build_title_slide(self, slide_data, idx, total):
        slide = self.prs.slides.add_slide(self._slide_layout)
        s = self.style_config
        _add_bg(slide, s["slide_bg"])

        stripe = _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 1.8, s["accent"])
        title = slide_data.get("title", "论文阅读汇报")
        _add_textbox(slide, 0.8, 0.4, 11.7, 1.2, title, Pt(40), "FFFFFF", True, PP_ALIGN.LEFT, s["font"])

        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            _add_textbox(slide, 0.8, 2.2, 11.7, 0.6, subtitle, Pt(20), s["muted"], False, PP_ALIGN.LEFT, s["font"])

        bullets = slide_data.get("bullets", [])
        _add_bullet_textbox(slide, 0.8, 3.0, 11.7, 3.0, [f"• {b}" for b in bullets], Pt(16), s["text"], s["font"])

        items = slide_data.get("visual_items", [])
        if items:
            _add_textbox(slide, 0.8, 6.0, 11.7, 0.5, "  →  ".join(items[:4]), Pt(13), s["accent2"], False, PP_ALIGN.LEFT, s["font"])

        _add_textbox(slide, 0.8, 6.8, 11.7, 0.4, f"AI 论文阅读生成 PPT  |  {self._slide_number(idx, total)}", Pt(10), s["muted"], False, PP_ALIGN.LEFT, s["font"])
        return slide

    def _build_content_slide(self, slide_data, idx, total):
        slide = self.prs.slides.add_slide(self._slide_layout)
        s = self.style_config
        bg = s["alt_bg"] if idx % 2 == 0 else s["slide_bg"]
        _add_bg(slide, bg)

        accent = s["accent2"] if idx % 3 == 0 else s["accent"]

        header_bar = _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.06, accent)

        title = slide_data.get("title", f"第 {idx} 页")
        _add_textbox(slide, 0.6, 0.25, 10, 0.65, title, s["title_size"], s["text"], True, PP_ALIGN.LEFT, s["font"])

        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            _add_textbox(slide, 0.6, 0.85, 10, 0.4, subtitle, s["subtitle_size"], s["muted"], False, PP_ALIGN.LEFT, s["font"])

        separator = _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.6, 1.25, 3.5, 0.03, accent)

        bullet_lines = [f"• {b}" for b in slide_data.get("bullets", [])[:5]]
        _add_bullet_textbox(slide, 0.6, 1.45, 5.5, 3.0, bullet_lines, s["body_size"], s["text"], s["font"])

        visual_items = slide_data.get("visual_items", ["问题", "方法", "实验", "结论"])
        visual_type = slide_data.get("visual_type", "flow")
        self._add_visual(slide, visual_type, visual_items, accent)

        zh_summary = slide_data.get("zh_summary", "")
        if zh_summary:
            summary_box = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 5.0, 6.0, 0.7, s["light_bg"])
            summary_box.text_frame.word_wrap = True
            summary_box.text_frame.paragraphs[0].text = str(zh_summary)[:200]
            summary_box.text_frame.paragraphs[0].font.size = s["small_size"]
            summary_box.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(s["text"])
            if s["font"]:
                summary_box.text_frame.paragraphs[0].font.name = s["font"]

        english_terms = slide_data.get("english_terms", [])
        if english_terms:
            term_lines = [str(t)[:120] for t in english_terms[:2]]
            term_box = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 5.85, 6.0, 0.65, s["highlight_bg"])
            tf = term_box.text_frame
            tf.word_wrap = True
            for i, line in enumerate(term_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = s["small_size"]
                p.font.color.rgb = _hex_to_rgb(s["muted"])
                if s["font"]:
                    p.font.name = s["font"]

        _add_textbox(slide, 0.6, 6.8, 11.7, 0.4, f"AI 论文阅读生成 PPT  |  {self._slide_number(idx, total)}", Pt(9), s["muted"], False, PP_ALIGN.LEFT, s["font"])

        slide_num_badge = _add_shape(slide, MSO_SHAPE.OVAL, 12.2, 0.2, 0.7, 0.7, accent)
        slide_num_badge.text_frame.paragraphs[0].text = str(idx)
        slide_num_badge.text_frame.paragraphs[0].font.size = Pt(14)
        slide_num_badge.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb("FFFFFF")
        slide_num_badge.text_frame.paragraphs[0].font.bold = True
        slide_num_badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return slide

    def _add_visual(self, slide, visual_type, items, accent):
        s = self.style_config
        items = items[:4]
        while len(items) < 2:
            items.append("维度")

        if visual_type == "flow":
            self._add_flow_visual(slide, items, accent, s)
        elif visual_type == "bars":
            self._add_bars_visual(slide, items, accent, s)
        elif visual_type == "matrix":
            self._add_matrix_visual(slide, items, accent, s)
        elif visual_type == "timeline":
            self._add_timeline_visual(slide, items, accent, s)
        else:
            self._add_flow_visual(slide, items, accent, s)

    def _add_flow_visual(self, slide, items, accent, s):
        _add_textbox(slide, 7.2, 1.45, 5.5, 0.4, "论文逻辑流程图", Pt(14), s["accent2"], True, PP_ALIGN.LEFT, s["font"])
        colors = [s["accent"], s["accent2"], s["accent"], s["accent2"]]
        for i, item in enumerate(items):
            x = 7.2 + i * 1.5
            y = 2.1
            _add_rounded_rect(slide, x, y, 1.3, 0.7, colors[i], f"{i+1}. {item}", Pt(11), "FFFFFF", s["font"])
            if i < len(items) - 1:
                _add_arrow(slide, x + 1.3, y + 0.25, 0.2, 0.15, s["muted"])
        _add_textbox(slide, 7.2, 3.1, 5.5, 0.8,
                     "用图把论文从问题、方法、证据到结论的阅读路径串起来。",
                     Pt(11), s["muted"], False, PP_ALIGN.LEFT, s["font"])

    def _add_bars_visual(self, slide, items, accent, s):
        _add_textbox(slide, 7.2, 1.45, 5.5, 0.4, "证据权重图", Pt(14), s["accent2"], True, PP_ALIGN.LEFT, s["font"])
        bar_widths = [4.5, 3.8, 3.0, 2.2]
        for i, item in enumerate(items):
            y = 2.1 + i * 0.65
            _add_textbox(slide, 7.2, y, 2.0, 0.35, item, Pt(11), s["text"], False, PP_ALIGN.LEFT, s["font"])
            track = _add_shape(slide, MSO_SHAPE.RECTANGLE, 9.3, y + 0.05, 4.0, 0.22, "E8E8E8")
            val_color = s["accent2"] if i % 2 else s["accent"]
            _add_shape(slide, MSO_SHAPE.RECTANGLE, 9.3, y + 0.05, bar_widths[i], 0.22, val_color)
        _add_textbox(slide, 7.2, 5.0, 5.5, 0.6,
                     "无明确数值时表示相对证据强弱，不虚构实验指标。",
                     Pt(11), s["muted"], False, PP_ALIGN.LEFT, s["font"])

    def _add_matrix_visual(self, slide, items, accent, s):
        _add_textbox(slide, 7.2, 1.45, 5.5, 0.4, "贡献对照矩阵", Pt(14), s["accent2"], True, PP_ALIGN.LEFT, s["font"])
        row_labels = items[:2] if len(items) >= 2 else ["现有问题", "本文方案"]
        col_labels = items[2:4] if len(items) >= 4 else ["技术贡献", "实验证据"]
        while len(row_labels) < 2:
            row_labels.append("维度")
        while len(col_labels) < 2:
            col_labels.append("维度")

        x0, y0 = 7.2, 2.1
        _add_rounded_rect(slide, x0 + 2.2, y0, 1.8, 0.5, s["accent"], col_labels[0], Pt(10), "FFFFFF", s["font"])
        _add_rounded_rect(slide, x0 + 4.2, y0, 1.8, 0.5, s["accent2"], col_labels[1], Pt(10), "FFFFFF", s["font"])
        for r in range(2):
            ry = y0 + 0.6 + r * 0.9
            cell_bg = "F0F0F0" if r % 2 == 0 else "FFFFFF"
            _add_rounded_rect(slide, x0, ry, 2.0, 0.7, "F5F5F5", row_labels[r], Pt(10), s["text"], s["font"])
            _add_rounded_rect(slide, x0 + 2.2, ry, 1.8, 0.7, cell_bg, "✓ 论文证据", Pt(10), s["text"], s["font"])
            _add_rounded_rect(slide, x0 + 4.2, ry, 1.8, 0.7, cell_bg, "→ 汇报重点", Pt(10), s["text"], s["font"])

    def _add_timeline_visual(self, slide, items, accent, s):
        _add_textbox(slide, 7.2, 1.45, 5.5, 0.4, "阅读路线图", Pt(14), s["accent2"], True, PP_ALIGN.LEFT, s["font"])
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 7.5, 2.85, 5.1, 0.05, s["accent"])
        for i, item in enumerate(items):
            x = 7.4 + i * 1.7
            circle = _add_shape(slide, MSO_SHAPE.OVAL, x, 2.5, 0.5, 0.5, s["accent2"] if i % 2 else s["accent"])
            circle.text_frame.paragraphs[0].text = str(i + 1)
            circle.text_frame.paragraphs[0].font.size = Pt(12)
            circle.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb("FFFFFF")
            circle.text_frame.paragraphs[0].font.bold = True
            circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            _add_textbox(slide, x - 0.2, 3.2, 1.5, 0.6, item, Pt(10), s["text"], False, PP_ALIGN.CENTER, s["font"])

    def add_slide(self, slide_data, idx, total):
        if idx == 1:
            return self._build_title_slide(slide_data, idx, total)
        return self._build_content_slide(slide_data, idx, total)

    def build(self, outline, output_path):
        slides_data = outline.get("slides", [])
        if not slides_data:
            slides_data = [{"title": "论文阅读汇报", "bullets": ["无可用内容"]}]
        total = len(slides_data)
        for i, slide_data in enumerate(slides_data, start=1):
            self.add_slide(slide_data, i, total)
            images = slide_data.get("images", [])
            if images:
                self._insert_images(self.prs.slides[-1], images)
        self.prs.save(output_path)
        logger.info(f"PPT saved to {output_path}")
        return output_path

    def _insert_images(self, slide, images):
        for img in images[:4]:
            try:
                data = img.get("data")
                if not data:
                    continue
                from io import BytesIO
                stream = BytesIO(data)
                w = img.get("w") or 400
                h = img.get("h") or 300
                slide.shapes.add_picture(
                    stream,
                    Inches(7.2), Inches(2.0),
                    Inches(min(w / 96, 5.5)),
                    Inches(min(h / 96, 4.0)),
                )
            except Exception as e:
                logger.warning(f"Failed to insert image: {e}")

    @staticmethod
    def import_template(template_path: str | Path) -> dict:
        """Analyze an existing PPTX and extract style information.

        Returns a style config dict compatible with STYLES.
        """
        from pptx import Presentation as PptxPresentation
        from pptx.util import Pt

        try:
            prs = PptxPresentation(str(template_path))
        except Exception as e:
            logger.warning(f"Cannot open template: {e}")
            return {}

        colors = set()
        fonts = set()
        bg_colors = set()

        for slide in prs.slides:
            try:
                bg = slide.background.fill
                if bg.type is not None:
                    bg_colors.add(str(bg.fore_color.rgb))
            except Exception:
                pass
            try:
                for shape in slide.shapes:
                    if hasattr(shape, 'fill'):
                        sf = shape.fill
                        if sf.type is not None:
                            bg_colors.add(str(sf.fore_color.rgb))
            except Exception:
                pass

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    colors.add(str(run.font.color.rgb))
                            except Exception:
                                pass
                            if run.font.name and run.font.name.strip():
                                fonts.add(run.font.name)
                        try:
                            if para.font.color and para.font.color.type is not None:
                                colors.add(str(para.font.color.rgb))
                        except Exception:
                            pass
                        if para.font.name and para.font.name.strip():
                            fonts.add(para.font.name)

        if not colors and not fonts:
            return {
                "name": "导入模板（默认）",
                "slide_bg": "FFFFFF",
                "alt_bg": "F2F2F2",
                "accent": "2050A0",
                "accent2": "0A7A6E",
                "text": "1C2833",
                "muted": "5D6D7E",
                "light_bg": "EBF5FB",
                "highlight_bg": "FEF9E7",
                "title_size": Pt(32),
                "subtitle_size": Pt(18),
                "body_size": Pt(14),
                "small_size": Pt(11),
                "font": "Microsoft YaHei",
            }

        color_list = sorted(colors)
        accent = color_list[len(color_list) // 2] if color_list else "2050A0"
        accent2 = color_list[0] if color_list else "2E86C1"
        bg_hex = "FFFFFF"
        if bg_colors:
            bg_hex = list(bg_colors)[0]
        font_name = list(fonts)[0] if fonts else "Microsoft YaHei"

        return {
            "name": "导入模板",
            "slide_bg": bg_hex,
            "alt_bg": _darken_hex(bg_hex, 10),
            "accent": accent,
            "accent2": accent2,
            "text": "1C2833",
            "muted": "5D6D7E",
            "light_bg": _lighten_hex(accent, 85),
            "highlight_bg": "FEF9E7",
            "title_size": Pt(32),
            "subtitle_size": Pt(18),
            "body_size": Pt(14),
            "small_size": Pt(11),
            "font": font_name,
        }


def _lighten_hex(hex_color, percent):
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    r = min(255, r + int((255 - r) * percent / 100))
    g = min(255, g + int((255 - g) * percent / 100))
    b = min(255, b + int((255 - b) * percent / 100))
    return f"{r:02X}{g:02X}{b:02X}"


def _darken_hex(hex_color, percent):
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    r = max(0, r - int(r * percent / 100))
    g = max(0, g - int(g * percent / 100))
    b = max(0, b - int(b * percent / 100))
    return f"{r:02X}{g:02X}{b:02X}"

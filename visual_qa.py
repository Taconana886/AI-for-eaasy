import logging
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class VisualQAReport:
    def __init__(self):
        self.issues: list[dict] = []
        self.passed = True

    def add_issue(self, slide_idx: int, severity: str, category: str, message: str):
        self.issues.append({
            "slide": slide_idx,
            "severity": severity,
            "category": category,
            "message": message,
        })
        if severity == "error":
            self.passed = False

    def summary(self) -> str:
        if not self.issues:
            return "视觉检查通过，无问题"
        by_severity = {}
        for issue in self.issues:
            by_severity.setdefault(issue["severity"], []).append(issue)
        parts = []
        for sev in ["error", "warning", "info"]:
            items = by_severity.get(sev, [])
            if items:
                parts.append(f"{sev.upper()}: {len(items)} 个")
        return f"视觉检查完成 ({', '.join(parts)})"


def check_pptx(pptx_path: str) -> VisualQAReport:
    """Check generated PPTX for visual quality issues.

    Checks:
    - Text overflow (text exceeds shape bounds)
    - Element overlap
    - Color contrast (light-on-light, dark-on-dark)
    - Empty slides
    """
    report = VisualQAReport()
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        report.add_issue(0, "error", "file", f"无法打开 PPTX: {e}")
        return report

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for slide_idx, slide in enumerate(prs.slides, start=1):
        shapes_info = []
        for shape in slide.shapes:
            s = {
                "name": shape.name or "",
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "has_text": shape.has_text_frame,
                "text": shape.text_frame.text if shape.has_text_frame else "",
            }
            shapes_info.append(s)

            if s["has_text"] and s["text"]:
                _check_text_overflow(shape, slide, slide_idx, report, s)
                _check_contrast(shape, slide, slide_idx, report, s)

        _check_overlap(shapes_info, slide_idx, report)
        _check_empty_slide(shapes_info, slide_idx, report)
        _check_out_of_bounds(shapes_info, slide_width, slide_height, slide_idx, report)

    logger.info(report.summary())
    return report


def _check_text_overflow(shape, slide, slide_idx, report, s):
    tf = shape.text_frame
    estimated_lines = 0
    for para in tf.paragraphs:
        text_len = len(para.text)
        if text_len == 0:
            continue
        chars_per_line = max(1, (shape.width or Inches(4)) // Pt(7))
        estimated_lines += max(1, text_len // chars_per_line + 1)

    max_lines = (shape.height or Inches(1)) // Pt(14)
    if max_lines <= 0:
        max_lines = 1

    if estimated_lines > max_lines * 1.5 and estimated_lines > 3:
        report.add_issue(
            slide_idx, "warning", "overflow",
            f"文本框「{s['name']}」文字可能溢出: 估算 {estimated_lines} 行, 空间约 {max_lines} 行"
        )


def _check_contrast(shape, slide, slide_idx, report, s):
    try:
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                font_color = run.font.color
                if font_color and font_color.rgb:
                    text_rgb = font_color.rgb
                    bg_rgb = _get_bg_color(slide)
                    if bg_rgb and _contrast_ratio(text_rgb, bg_rgb) < 2.5:
                        report.add_issue(
                            slide_idx, "warning", "contrast",
                            f"文字颜色 {text_rgb} 与背景 {bg_rgb} 对比度不足"
                        )
                        return
    except Exception:
        pass


def _get_bg_color(slide):
    try:
        bg = slide.background.fill
        if bg.type is not None:
            return bg.fore_color.rgb
    except Exception:
        pass
    return RGBColor(0xFF, 0xFF, 0xFF)


def _relative_luminance(rgb):
    def linearize(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * linearize(rgb[0]) + 0.7152 * linearize(rgb[1]) + 0.0722 * linearize(rgb[2])


def _contrast_ratio(rgb1, rgb2):
    l1 = _relative_luminance(rgb1)
    l2 = _relative_luminance(rgb2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _check_overlap(shapes, slide_idx, report):
    for i, a in enumerate(shapes):
        for j, b in enumerate(shapes):
            if i >= j:
                continue
            if not a["has_text"] or not b["has_text"]:
                continue
            if not a["text"] or not b["text"]:
                continue
            ax1, ay1 = a["left"], a["top"]
            ax2, ay2 = ax1 + a["width"], ay1 + a["height"]
            bx1, by1 = b["left"], b["top"]
            bx2, by2 = bx1 + b["width"], by1 + b["height"]

            overlap_x = max(0, min(ax2, bx2) - max(ax1, bx1))
            overlap_y = max(0, min(ay2, by2) - max(ay1, by1))

            if overlap_x > 0 and overlap_y > 0:
                overlap_area = overlap_x * overlap_y
                a_area = a["width"] * a["height"]
                if a_area > 0 and overlap_area / a_area > 0.3:
                    report.add_issue(
                        slide_idx, "info", "overlap",
                        f"「{a['name']}」与「{b['name']}」可能重叠"
                    )


def _check_empty_slide(shapes, slide_idx, report):
    non_empty = [s for s in shapes if s["has_text"] and s["text"].strip()]
    if len(non_empty) <= 1:
        report.add_issue(slide_idx, "warning", "empty", "幻灯片内容过少")


def _check_out_of_bounds(shapes, slide_w, slide_h, slide_idx, report):
    margin = Inches(0.2)
    for s in shapes:
        if not s["has_text"] or not s["text"]:
            continue
        if s["left"] < -margin or s["top"] < -margin:
            report.add_issue(slide_idx, "info", "bounds", f"「{s['name']}」超出幻灯片边界")


def auto_fix_pptx(pptx_path: str) -> str:
    """Attempt automatic fixes for common issues.

    Returns path to fixed file.
    """
    from pathlib import Path
    import tempfile
    import shutil

    report = check_pptx(pptx_path)
    if report.passed:
        return pptx_path

    fixed_path = str(Path(pptx_path).with_suffix(".fixed.pptx"))
    shutil.copy2(pptx_path, fixed_path)
    try:
        prs = Presentation(fixed_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    _shrink_text_to_fit(shape)
        prs.save(fixed_path)
        logger.info(f"Auto-fixed PPTX saved to {fixed_path}")
        return fixed_path
    except Exception as e:
        logger.warning(f"Auto-fix failed: {e}")
        return pptx_path


def _shrink_text_to_fit(shape):
    try:
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size > Pt(40):
                    run.font.size = Pt(36)
    except Exception:
        pass
